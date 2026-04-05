# -*- coding: utf-8 -*-
"""
Create improved TSRS Presentation using python-pptx
7 slides with dark professional theme, diagrams, and references.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette (Dark Navy + Teal) ===
NAVY = RGBColor(0x0F, 0x17, 0x2A)
DARK_BLUE = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
LIGHT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
VERY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, x, y, w, h, size=14, color=WHITE, bold=False, align=PP_ALIGN.RIGHT, font='Arial'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_rich_text(slide, lines, x, y, w, h, align=PP_ALIGN.RIGHT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# =====================================================
# SLIDE 1: Title
# =====================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1, NAVY)

# Accent line top
add_rect(slide1, 0, 0, 13.333, 0.06, TEAL)

# Shield icon area
add_rounded_rect(slide1, 5.5, 1.2, 2.3, 2.3, DARK_BLUE, TEAL)
add_text(slide1, '🛡️', 5.8, 1.3, 1.7, 1.8, size=60, align=PP_ALIGN.CENTER)

# Title
add_text(slide1, 'TSRS', 1.5, 3.8, 10.3, 1.0, size=54, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide1, 'Tsunami Station Risk Score', 1.5, 4.7, 10.3, 0.6, size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide1, 'מודל גאוגרפי לתיעדוף הקצאת משאבים — משטרת ישראל', 1.5, 5.3, 10.3, 0.5, size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide1, 0, 6.9, 13.333, 0.6, DARK_BLUE)
add_text(slide1, 'מבוסס על: Wood & Jones (2015) | Tucson Workload Model | מחקרי ABM | יוני 2025', 0.5, 6.95, 12.3, 0.4, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 2: Problem & Vision
# =====================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, NAVY)
add_rect(slide2, 0, 0, 13.333, 0.06, TEAL)

add_text(slide2, 'הבעיה: ניהול אינטואיטיבי בסביבה קריטית', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# Left column - Problems
add_rounded_rect(slide2, 0.5, 1.3, 5.8, 4.5, DARK_BLUE, RGBColor(0x33, 0x41, 0x55))

problems = [
    ('פער מידע', 'אין מפת הצפות ייעודית המשולבת עם טריטוריות תחנות משטרה'),
    ('הקצאה אינטואיטיבית', 'כוחות מוקצים ידנית ללא שיקלול פגיעות אוכלוסייה, צפיפות ותנועה'),
    ('זמן התראה קצר', 'ים תיכון = Near-field. פחות מ-20 דקות לתגובה. אין מקום לאלתור'),
    ('משימות המשטרה', 'פינוי אוכלוסייה, ניהול תנועה, סגירת זירות ומניעת ביזה'),
]

y_pos = 1.5
for title, desc in problems:
    add_text(slide2, f'🔴  {title}', 0.8, y_pos, 5.2, 0.35, size=14, color=WHITE, bold=True)
    add_text(slide2, desc, 1.2, y_pos + 0.35, 4.8, 0.55, size=12, color=LIGHT_GRAY)
    y_pos += 1.0

# Right column - Vision
add_rounded_rect(slide2, 6.8, 1.3, 6.0, 4.5, DARK_BLUE, TEAL)
add_text(slide2, '🎯  החזון', 6.9, 1.5, 5.7, 0.4, size=20, color=TEAL, bold=True)
add_text(slide2, 'מפקד תחנה פותח מסך, בוחר גובה גל ונפה, ורואה:', 7.1, 2.0, 5.3, 0.4, size=14, color=WHITE)

vision_items = [
    'מפת סיכון של טריטוריית התחנה',
    'ציון TSRS מחושב (0-100) לכל עיר',
    'מספר ניידות גיבוי נדרש',
    'צירי פינוי מומלצים ומקלטים אנכיים',
    'הנחיות מבצעיות מותאמות לתרחיש',
]
y_v = 2.5
for item in vision_items:
    add_text(slide2, f'✅  {item}', 7.1, y_v, 5.3, 0.35, size=13, color=VERY_LIGHT)
    y_v += 0.4

# Assumptions box
add_rounded_rect(slide2, 0.5, 6.0, 12.3, 1.2, RGBColor(0x15, 0x20, 0x35), TEAL)
add_text(slide2, 'הנחות בסיס (כולל הנחות סמויות)', 0.8, 6.05, 11.7, 0.35, size=13, color=TEAL, bold=True)
assumptions = '• המשטרה היא First Responder — פועלת לפני פיקוד העורף ומל"ח  •  נתוני DEM ואוכלוסייה זמינים מלמ"ס ומפ"י  •  מודל הצפה מבוסס על גובה גל בלבד (לא בתימות/מקור רעידה)  •  ABM מבוסס על הנחות תנועה ממוצעות  •  אשכול סוציו-אקונומי כ-proxy לניידות ופגיעות'
add_text(slide2, assumptions, 0.8, 6.4, 11.7, 0.7, size=11, color=LIGHT_GRAY)

# =====================================================
# SLIDE 3: Model - 5 Components
# =====================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, NAVY)
add_rect(slide3, 0, 0, 13.333, 0.06, TEAL)

add_text(slide3, 'המודל: 5 מדדים משוקללים', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# Formula box
add_rounded_rect(slide3, 1.0, 1.2, 11.3, 0.7, DARK_BLUE, TEAL)
add_text(slide3, 'TSRS = (H × 0.35) + (V × 0.30) + (O × 0.20) + (R⁻¹ × 0.10) + (I⁻¹ × 0.05)', 1.2, 1.25, 10.9, 0.6, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font='Courier New')

# 5 component cards
components = [
    ('H', '35%', 'חשיפה פיזית', 'קרבת חוף, % הצפה, זמן גל', RED),
    ('V', '30%', 'פגיעות אוכלוסייה', 'גיל, ניידות, מוסדות, מדד ח"כ', RGBColor(0x8B, 0x5C, 0xF6)),
    ('O', '20%', 'עומס מבצעי', 'ABM, צווארי בקבוק, צירי פינוי', ORANGE),
    ('R⁻¹', '10%', 'גרעון משאבים', 'ניידות, זמן תגובה, KYT', GREEN),
    ('I⁻¹', '5%', 'חוסן תשתית', 'צירים חלופיים, מקלטים, DTN', BLUE),
]

x_pos = 0.5
for key, weight, title, desc, color in components:
    card = add_rounded_rect(slide3, x_pos, 2.2, 2.35, 2.8, DARK_BLUE, color)
    add_text(slide3, key, x_pos + 0.1, 2.3, 2.15, 0.5, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
    # Weight badge
    badge = add_rounded_rect(slide3, x_pos + 0.7, 2.9, 0.95, 0.35, color)
    add_text(slide3, weight, x_pos + 0.7, 2.9, 0.95, 0.35, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, title, x_pos + 0.1, 3.4, 2.15, 0.4, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide3, desc, x_pos + 0.1, 3.8, 2.15, 0.8, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    x_pos += 2.5

# Risk tiers
add_text(slide3, 'דרגות סיכון:', 0.8, 5.3, 2.0, 0.4, size=14, color=WHITE, bold=True)
tiers = [('80+', 'קריטי', RED), ('60+', 'גבוה', ORANGE), ('40+', 'בינוני', YELLOW), ('20+', 'נמוך', GREEN), ('<20', 'מינימלי', BLUE)]
tx = 3.0
for score, label, color in tiers:
    badge = add_rounded_rect(slide3, tx, 5.3, 1.8, 0.4, color)
    add_text(slide3, f'{score}  {label}', tx, 5.3, 1.8, 0.4, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    tx += 2.0

# =====================================================
# SLIDE 4: Process Diagram (INPUT → PROCESSING → OUTPUT)
# =====================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, NAVY)
add_rect(slide4, 0, 0, 13.333, 0.06, TEAL)

add_text(slide4, 'תרשים המודל: קלט → עיבוד → פלט', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# INPUT column
add_rounded_rect(slide4, 0.3, 1.3, 3.8, 5.5, DARK_BLUE, BLUE)
add_text(slide4, '📥  נתוני קלט', 0.5, 1.4, 3.4, 0.5, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

inputs = [
    'DEM 5m — מפ"י (מודל גובה דיגיטלי)',
    'גבולות מוניציפליים — OSM',
    'אוכלוסייה וגילאים — למ"ס',
    'אשכול סוציו-אקונומי — למ"ס',
    'רשת כבישים — OSM',
    'תחנות משטרה — OSM',
    'מודל הצפה — GSI (גובה גל)',
]
yi = 2.1
for inp in inputs:
    add_text(slide4, f'▸  {inp}', 0.5, yi, 3.4, 0.35, size=11, color=VERY_LIGHT)
    yi += 0.42

# ARROWS
add_text(slide4, '➡', 4.2, 3.5, 0.5, 0.5, size=30, color=TEAL, align=PP_ALIGN.CENTER)

# PROCESSING column
add_rounded_rect(slide4, 4.8, 1.3, 3.7, 5.5, DARK_BLUE, TEAL)
add_text(slide4, '⚙️  שלבי עיבוד', 5.0, 1.4, 3.3, 0.5, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

steps = [
    ('1', 'חישוב אזורי הצפה לפי גובה גל'),
    ('2', 'חיתוך מרחבי עם גבולות ערים'),
    ('3', 'חישוב H — אחוז שטח מוצף'),
    ('4', 'חישוב V — פגיעות דמוגרפית מלמ"ס'),
    ('5', 'חישוב O — סימולציית ABM לפינוי'),
    ('6', 'חישוב R⁻¹ — הערכת משאבי תגובה'),
    ('7', 'חישוב I⁻¹ — זיהוי מקלטים ותשתיות'),
    ('8', 'נרמול Min-Max + שקלול TSRS'),
]
yp = 2.1
for num, step in steps:
    add_text(slide4, f'{num}.  {step}', 5.0, yp, 3.3, 0.35, size=11, color=VERY_LIGHT)
    yp += 0.42

# ARROWS
add_text(slide4, '➡', 8.6, 3.5, 0.5, 0.5, size=30, color=TEAL, align=PP_ALIGN.CENTER)

# OUTPUT column
add_rounded_rect(slide4, 9.2, 1.3, 3.8, 5.5, DARK_BLUE, GREEN)
add_text(slide4, '📤  תוצרי פלט', 9.4, 1.4, 3.4, 0.5, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Geographic outputs
add_text(slide4, 'תוצרים גאוגרפיים:', 9.4, 2.1, 3.4, 0.3, size=13, color=WHITE, bold=True)
geo_outputs = ['מפת חום הצפות דינמית', 'פוליגוני ערים צבועים לפי TSRS', 'נקודות תחנות משטרה', 'שכבות OSM (כבישים, מבנים)']
yg = 2.5
for out in geo_outputs:
    add_text(slide4, f'🗺️  {out}', 9.4, yg, 3.4, 0.3, size=11, color=VERY_LIGHT)
    yg += 0.35

# Tabular outputs
add_text(slide4, 'תוצרים טבלאיים:', 9.4, 4.1, 3.4, 0.3, size=13, color=WHITE, bold=True)
tab_outputs = ['ציון TSRS (0-100) לכל עיר', 'פירוט 5 מדדים (H,V,O,R,I)', 'הנחיות מבצעיות בעברית', 'פילוח דמוגרפי (גילאים, ח"כ)', 'ייצוא PDF לתחנה']
yt = 4.5
for out in tab_outputs:
    add_text(slide4, f'📊  {out}', 9.4, yt, 3.4, 0.3, size=11, color=VERY_LIGHT)
    yt += 0.35

# =====================================================
# SLIDE 5: Application Screenshots Description
# =====================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, NAVY)
add_rect(slide5, 0, 0, 13.333, 0.06, TEAL)

add_text(slide5, 'היישום: GIS אינטראקטיבי לתיעדוף בזמן אמת', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# Feature cards (2x3 grid)
features = [
    ('🗺️', 'מפת GIS', 'מפת Leaflet עם שכבות OSM, ESRI\nHillshade, ומפת חום הצפות דינמית'),
    ('🏙️', 'גבולות ערים', '554 גבולות מוניציפליים אמיתיים מ-OSM\nצבועים לפי ציון TSRS'),
    ('🌊', 'סליידר גובה גל', 'בחירת גובה גל 0.5-10 מ\' עם\nעדכון מיידי של מפת החום'),
    ('📊', 'פילוח דמוגרפי', 'התפלגות גילאים, אשכול ח"כ\nמנתוני למ"ס אמיתיים'),
    ('🛡️', 'תחנות משטרה', 'נתוני OSM אמיתיים עם אייקון\nמשטרת ישראל על המפה'),
    ('📋', 'הנחיות מבצעיות', 'ניידות גיבוי, צירי פינוי,\nמקלטים אנכיים — לכל עיר'),
]

positions = [(0.5, 1.3), (4.6, 1.3), (8.7, 1.3), (0.5, 4.2), (4.6, 4.2), (8.7, 4.2)]
for (icon, title, desc), (fx, fy) in zip(features, positions):
    add_rounded_rect(slide5, fx, fy, 3.8, 2.5, DARK_BLUE, RGBColor(0x33, 0x41, 0x55))
    add_text(slide5, icon, fx + 0.2, fy + 0.2, 0.6, 0.6, size=28, align=PP_ALIGN.CENTER)
    add_text(slide5, title, fx + 0.8, fy + 0.2, 2.7, 0.4, size=16, color=WHITE, bold=True)
    add_text(slide5, desc, fx + 0.3, fy + 0.8, 3.2, 1.4, size=12, color=LIGHT_GRAY)

# =====================================================
# SLIDE 6: Data Sources & References
# =====================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, NAVY)
add_rect(slide6, 0, 0, 13.333, 0.06, TEAL)

add_text(slide6, 'מקורות מידע וביבליוגרפיה', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# Data sources (left)
add_rounded_rect(slide6, 0.5, 1.3, 6.0, 3.5, DARK_BLUE, RGBColor(0x33, 0x41, 0x55))
add_text(slide6, '📂  מקורות נתונים', 0.7, 1.4, 5.6, 0.4, size=16, color=WHITE, bold=True)

sources = [
    'הלמ"ס — נתוני אוכלוסייה, גילאים, אשכול סוציו-אקונומי',
    'מפ"י — מודל גובה דיגיטלי (DEM) 5 מטר',
    'GSI (מכון גיאולוגי) — מודל הצפות צונאמי',
    'OpenStreetMap — גבולות מוניציפליים, כבישים, מבנים, תחנות',
    'ESRI — שכבת Hillshade (תבליט)',
    'GOVMAP — שכבות WMS ממשלתיות',
]
ys = 2.0
for src in sources:
    add_text(slide6, f'▸  {src}', 0.7, ys, 5.6, 0.35, size=12, color=VERY_LIGHT)
    ys += 0.42

# References (right)
add_rounded_rect(slide6, 6.8, 1.3, 6.0, 3.5, DARK_BLUE, TEAL)
add_text(slide6, '📚  מקורות ביבליוגרפיים', 7.0, 1.4, 5.6, 0.4, size=16, color=TEAL, bold=True)

refs = [
    '1. Wood, N. & Jones, J. (2015). "Tsunami Exposure\n   Assessment" — USGS Scientific Report.',
    '2. Tucson Police Dept. — Workload-Based Staffing\n   Model for Resource Allocation.',
    '3. Mas, E. et al. (2012). "Agent-Based Simulation\n   of Tsunami Evacuation" — Natural Hazards.',
    '4. Patel, D. (2019). "GIS-Based Multi-Criteria\n   Tsunami Risk Assessment" — Applied Geography.',
    '5. Israel CBS (2022). "Socio-Economic Profiles\n   of Local Authorities" — Statistical Abstract.',
    '6. Geist, E. & Parsons, T. (2006). "Probabilistic\n   Analysis of Tsunami Hazards" — Natural Hazards.',
    '7. Leone, F. et al. (2011). "Evacuation Modeling\n   for Tsunami Risk" — Natural Hazards & Earth.',
]
yr = 1.9
for ref in refs:
    add_text(slide6, ref, 7.0, yr, 5.6, 0.5, size=10, color=VERY_LIGHT, font='Courier New')
    yr += 0.42

# Tech stack
add_rounded_rect(slide6, 0.5, 5.1, 12.3, 1.8, DARK_BLUE, RGBColor(0x33, 0x41, 0x55))
add_text(slide6, '🔧  סטאק טכנולוגי', 0.7, 5.2, 11.9, 0.4, size=16, color=WHITE, bold=True)

tech = [
    ('Frontend', 'Leaflet.js + HTML/CSS/JS + Firebase'),
    ('Backend', 'Python FastAPI + GeoJSON'),
    ('Map Services', 'OSM + ESRI Hillshade + GOVMAP WMS'),
    ('Data', 'CBS CSV + Overpass API + Leaflet.heat'),
]
xt = 0.7
for label, value in tech:
    add_text(slide6, label, xt, 5.7, 2.8, 0.3, size=12, color=TEAL, bold=True)
    add_text(slide6, value, xt, 6.0, 2.8, 0.3, size=11, color=LIGHT_GRAY)
    xt += 3.0

# =====================================================
# SLIDE 7: Summary & Call to Action
# =====================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, NAVY)
add_rect(slide7, 0, 0, 13.333, 0.06, TEAL)

add_text(slide7, 'סיכום וקריאה לפעולה', 0.8, 0.3, 11.7, 0.7, size=30, color=TEAL, bold=True)

# Summary cards (left)
summary_items = [
    ('5 מדדים', 'ציון TSRS 0-100 לכל עיר'),
    ('GIS + ABM + למ"ס', 'מודל מבוסס נתונים אמיתיים'),
    ('ממשק אינטואיטיבי', 'סליידר + נפה + מפת חום'),
    ('554 ערים', 'גבולות מוניציפליים מלאים'),
    ('5 שפות', 'עברית, אנגלית, רוסית, צרפתית, ספרדית'),
    ('Open Source', 'OSM + Leaflet — עלות מינימלית'),
]

ys = 1.2
for title, desc in summary_items:
    add_rounded_rect(slide7, 0.5, ys, 5.5, 0.55, DARK_BLUE, RGBColor(0x33, 0x41, 0x55))
    add_text(slide7, f'✅  {title}', 0.7, ys + 0.05, 2.3, 0.45, size=13, color=TEAL, bold=True)
    add_text(slide7, desc, 3.0, ys + 0.05, 2.8, 0.45, size=12, color=LIGHT_GRAY)
    ys += 0.65

# Call to action (right)
add_rounded_rect(slide7, 6.5, 1.2, 6.3, 4.8, DARK_BLUE, TEAL)
add_text(slide7, '🚀  הצעדים הבאים', 6.7, 1.3, 5.9, 0.5, size=20, color=TEAL, bold=True)

cta = [
    ('1', 'אישור BRD על ידי ממ"מ / מנהל מחשוב'),
    ('2', 'תיאום GSI ומפ"י לנתוני הצפה מדויקים'),
    ('3', 'בחירת 3 תחנות פיילוט בחוף הים התיכון'),
    ('4', 'הקמת סביבה (Docker + DB) — שבוע 1'),
    ('5', 'Kick-off: צוות GIS + מפקדי תחנות'),
]

yc = 2.0
for num, text in cta:
    add_rounded_rect(slide7, 6.8, yc, 0.5, 0.5, TEAL)
    add_text(slide7, num, 6.8, yc, 0.5, 0.5, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide7, text, 7.5, yc + 0.05, 5.0, 0.45, size=14, color=WHITE)
    yc += 0.7

# Quote
add_rect(slide7, 0.5, 6.3, 12.3, 0.8, DARK_BLUE)
add_text(slide7, '"מידע גאוגרפי מציל חיים — הכנה מוקדמת מציל יותר"', 0.5, 6.35, 12.3, 0.6, size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'TSRS_Presentation_Improved.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Slides: 7')
